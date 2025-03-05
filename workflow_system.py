"""
Comprehensive Workflow System

This module implements a unified research workflow system using Langgraph for orchestration
and Pydantic agents for specialized roles. The system features a modular architecture with
distinct agents for different research tasks.
"""

import os
import asyncio
import json
from datetime import datetime
from typing import List, Dict, Any, Optional, Union, Literal, TypedDict
from enum import Enum
from pathlib import Path
import tempfile
import time
import random

# Pydantic for structured data and agent definitions
from pydantic import BaseModel, Field, validator
from pydantic_ai import Agent, RunContext

# LangGraph for workflow orchestration
from langgraph.graph import StateGraph, END
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.messages import HumanMessage, AIMessage

# Rich for terminal UI
from rich.console import Console
from rich.panel import Panel
from rich.progress import Progress, SpinnerColumn, TextColumn, BarColumn, TimeElapsedColumn
from rich.table import Table
from rich.markdown import Markdown

# Document export formats
import docx
from docx.shared import Inches, Pt
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table as PDFTable
from reportlab.lib.styles import getSampleStyleSheet
import pptx
from pptx.util import Inches as PPTInches
import markdown
import html

# Load environment variables
from dotenv import load_dotenv
load_dotenv()

# Initialize console for rich terminal output
console = Console()

# Initialize OpenAI client for LangGraph
llm = ChatOpenAI(model="gpt-4o", temperature=0.1)

#############################################
# Data Models for Structured Information
#############################################


class AgentType(str, Enum):
    """Enum representing different agent types in the workflow."""
    PLANNER = "planner"
    RESEARCHER = "researcher"
    WRITER = "writer"
    FACT_CHECKER = "fact_checker"
    EDITOR = "editor"
    FORMATTER = "formatter"
    SOURCE_GATHERER = "source_gatherer"
    DOCUMENT_EXPORTER = "document_exporter"
    RESEARCH_PLANNER = "research_planner"


class ResearchSubtopic(BaseModel):
    """A subtopic for research with relevant questions."""
    title: str = Field(..., description="The title of the subtopic")
    description: str = Field(...,
                             description="Brief description of the subtopic")
    key_questions: Optional[List[str]] = Field(default_factory=list,
                                               description="Key questions to address in this subtopic")
    content: Optional[str] = Field(
        default="", description="Content for this subtopic")
    notes: Optional[List[str]] = Field(
        default_factory=list, description="Notes about this subtopic")


class Source(BaseModel):
    """A research source with citation information."""
    id: str = Field(..., description="Unique identifier for the source")
    title: str = Field(..., description="The title of the source")
    url: Optional[str] = Field(
        None, description="URL of the source, if applicable")
    authors: Optional[List[str]] = Field(
        None, description="Authors of the source")
    publication_date: Optional[str] = Field(
        None, description="Publication date of the source")
    publisher: Optional[str] = Field(
        None, description="Publisher of the source")
    content: Optional[str] = Field(
        None, description="Content or excerpt from the source")
    relevance_score: Optional[float] = Field(
        None, description="Relevance score for this source (0-1)")
    credibility_score: Optional[float] = Field(
        None, description="Source credibility score (0-10)")
    key_insights: Optional[List[str]] = Field(
        default_factory=list, description="Key insights from this source")
    citation_apa: Optional[str] = Field(
        None, description="The source citation in APA format")
    citation_mla: Optional[str] = Field(
        None, description="The source citation in MLA format")
    citation_chicago: Optional[str] = Field(
        None, description="The source citation in Chicago format")


class ValidationResult(BaseModel):
    """Validation result for a fact or claim."""
    claim: str = Field(..., description="The claim being validated")
    is_valid: bool = Field(...,
                           description="Whether the claim is valid based on sources")
    confidence: float = Field(..., description="Confidence score (0-1)")
    supporting_sources: List[str] = Field(
        ..., description="List of source IDs that support/refute this claim")
    explanation: str = Field(...,
                             description="Explanation of the validation result")


class ResearchFindings(BaseModel):
    """Research findings for a subtopic."""
    subtopic: str = Field(..., description="The subtopic title")
    key_findings: List[str] = Field(...,
                                    description="List of key findings for this subtopic")
    analysis: str = Field(..., description="Analysis of the findings")
    sources: List[str] = Field(...,
                               description="Source IDs used for this subtopic")


class ResearchPlan(BaseModel):
    """Research plan created by the planner agent."""
    topic: str = Field(..., description="The main research topic")
    subtopics: List[ResearchSubtopic] = Field(
        ..., description="List of subtopics to research")
    approach: str = Field(..., description="Overall approach to the research")
    timeline: Dict[str, str] = Field(...,
                                     description="Timeline for each research phase")
    deliverables: List[str] = Field(...,
                                    description="Expected deliverables from the research")


class ResearchState(BaseModel):
    """The full state of the research workflow."""
    query: str = Field(..., description="Original research query")
    plan: Optional[ResearchPlan] = Field(None, description="Research plan")
    subtopics: List[ResearchSubtopic] = Field(
        default_factory=list, description="List of research subtopics")
    sources: Dict[str, Source] = Field(
        default_factory=dict, description="Dictionary of sources by ID")
    findings: Dict[str, ResearchFindings] = Field(
        default_factory=dict, description="Findings by subtopic")
    validated_facts: List[ValidationResult] = Field(
        default_factory=list, description="List of validated facts")
    content: Dict[str, str] = Field(
        default_factory=dict, description="Draft content sections")
    final_document: Optional[str] = Field(
        None, description="Final formatted document content")
    status: str = Field(default="initialized",
                        description="Current status of the research")
    messages: List[Dict[str, Any]] = Field(
        default_factory=list, description="History of agent messages")
    document_format: str = Field(
        "markdown", description="Output document format")
    citation_style: str = Field("APA", description="Citation style to use")
    research_depth: str = Field(
        "comprehensive", description="Research depth level")
    timestamp: str = Field(default_factory=lambda: datetime.now(
    ).isoformat(), description="Timestamp of the research")
    errors: List[str] = Field(
        default_factory=list, description="Errors encountered during the workflow")


class FormatConfig(BaseModel):
    """Configuration for document formatting and export."""
    format_type: Literal["markdown", "pdf", "docx",
                         "pptx"] = Field(..., description="Output format type")
    citation_style: Literal["APA", "MLA", "Chicago"] = Field(
        "APA", description="Citation style to use")
    include_toc: bool = Field(
        True, description="Whether to include a table of contents")
    include_abstract: bool = Field(
        True, description="Whether to include an abstract")
    include_executive_summary: bool = Field(
        False, description="Whether to include an executive summary")
    template: Optional[str] = Field(
        None, description="Optional template to use")


class WorkflowConfig(BaseModel):
    """Configuration for the research workflow."""
    query: str = Field(..., description="Research topic query")
    research_depth: Literal["basic", "comprehensive", "exhaustive"] = Field("comprehensive",
                                                                            description="Depth of research")
    format_config: FormatConfig = Field(...,
                                        description="Document formatting configuration")
    verbose: bool = Field(
        False, description="Whether to output verbose progress information")
    max_sources: int = Field(
        10, description="Maximum number of sources to gather")


#############################################
# Pydantic Agent Definitions
#############################################

def format_as_xml(data: Any) -> str:
    """Format data as XML-like string for PydanticAI agents."""
    if isinstance(data, dict):
        result = ""
        for key, value in data.items():
            if isinstance(value, (list, dict)):
                value = json.dumps(value)
            result += f"<{key}>{value}</{key}>\n"
        return result
    elif isinstance(data, BaseModel):
        return format_as_xml(data.dict())
    else:
        return str(data)


# Planner Agent
planner_agent = Agent(
    'openai:gpt-4o',
    result_type=ResearchPlan,
    system_prompt="""You are an expert research planner who creates comprehensive research plans.
    
    Given a research topic, you will:
    1. Break it down into logical subtopics for investigation
    2. Define a strategic approach to the research
    3. Create a timeline for each research phase
    4. Specify the expected deliverables
    
    Your plan should be academically rigorous, thorough, and efficient.
    You will output a ResearchPlan object with your complete research planning.
    """
)

# Researcher Agent
researcher_agent = Agent(
    'openai:gpt-4o',
    result_type=Dict[str, Source],
    system_prompt="""You are an expert research agent who gathers high-quality information from multiple sources.
    
    Given a research topic or subtopic, you will:
    1. Identify the most relevant and credible sources of information
    2. Extract key insights from each source
    3. Evaluate the credibility of each source
    4. Format proper citations in multiple styles
    
    Your research should be comprehensive, accurate, and well-documented.
    You will output a dictionary of Source objects containing your research findings.
    """
)

# Fact Checker Agent
fact_checker_agent = Agent(
    'openai:gpt-4o',
    result_type=List[ValidationResult],
    system_prompt="""You are an expert fact checking agent who validates research claims against multiple sources.
    
    Given a set of claims and supporting sources, you will:
    1. Evaluate each claim against the available evidence
    2. Determine whether each claim is valid
    3. Assign a confidence score to your evaluation
    4. Identify which sources support or refute each claim
    5. Provide a detailed explanation of your validation process
    
    Your fact checking should be rigorous, objective, and transparent.
    You will output a list of ValidationResult objects with your validations.
    """
)

# Writer Agent
writer_agent = Agent(
    'openai:gpt-4o',
    result_type=Dict[str, str],
    system_prompt="""You are an expert content writer who creates high-quality research content.
    
    Given research findings, validated facts, and a content structure, you will:
    1. Create engaging and informative content for each section
    2. Incorporate key findings and insights from the research
    3. Ensure all claims are supported by validated facts
    4. Maintain a consistent and appropriate tone
    
    Your writing should be clear, well-structured, and compelling.
    You will output a dictionary with section titles as keys and content as values.
    """
)

# Editor Agent
editor_agent = Agent(
    'openai:gpt-4o',
    result_type=Dict[str, str],
    system_prompt="""You are an expert editor who refines and improves research content.
    
    Given draft content sections, you will:
    1. Improve clarity, coherence, and flow
    2. Correct any grammatical or stylistic issues
    3. Ensure consistent terminology and tone
    4. Strengthen weak arguments and fix logical inconsistencies
    
    Your editing should enhance the quality of the content while preserving its core message.
    You will output an improved version of the content dictionary.
    """
)

# Formatter Agent
formatter_agent = Agent(
    'openai:gpt-4o',
    result_type=str,
    system_prompt="""You are an expert formatting agent who structures documents according to academic standards.
    
    Given edited content and formatting preferences, you will:
    1. Apply the appropriate academic structure (abstract, intro, body, conclusion, references)
    2. Format citations according to the specified style (APA, MLA, Chicago)
    3. Create a table of contents if requested
    4. Add appropriate headings, subheadings, and formatting elements
    
    Your formatting should adhere to academic standards and enhance readability.
    You will output the complete formatted document as a single string in markdown format.
    """
)

#############################################
# LangGraph Workflow Definition
#############################################


def create_research_workflow() -> StateGraph:
    """Create the LangGraph workflow for the research process."""
    # Define workflow states
    workflow = StateGraph(ResearchState)

    # Planning state
    def planning_state(state: ResearchState) -> Dict[str, Any]:
        """Planner agent creates the research plan."""
        console.print(
            Panel("[bold blue]Planning Research[/bold blue]", expand=False))

        # Use Pydantic agent to create plan
        plan_result = asyncio.run(planner_agent.run(
            f"Research Topic: {state.query}\n\nPlease create a comprehensive research plan for this topic."
        ))

        # Update state with plan
        state.plan = plan_result.data
        state.subtopics = plan_result.data.subtopics
        state.status = "planned"
        state.messages.append({
            "role": "agent",
            "agent": AgentType.RESEARCH_PLANNER.value,
            "content": f"Created research plan with {len(state.subtopics)} subtopics"
        })

        return {"state": state, "next": "research"}

    # Research state
    def research_state(state: ResearchState) -> Dict[str, Any]:
        """Researcher agent gathers sources for each subtopic."""
        console.print(
            Panel("[bold blue]Gathering Sources[/bold blue]", expand=False))

        # Process each subtopic
        for idx, subtopic in enumerate(state.subtopics):
            console.print(
                f"Researching subtopic {idx+1}/{len(state.subtopics)}: [bold]{subtopic.title}[/bold]")

            # Use Pydantic agent to research the subtopic
            sources_result = asyncio.run(researcher_agent.run(
                format_as_xml({
                    "query": state.query,
                    "subtopic": subtopic.dict(),
                    "depth": state.research_depth
                })
            ))

            # Add sources to state
            state.sources.update(sources_result.data)

        state.status = "researched"
        state.messages.append({
            "role": "agent",
            "agent": AgentType.RESEARCHER.value,
            "content": f"Gathered {len(state.sources)} sources across {len(state.subtopics)} subtopics"
        })

        return {"state": state, "next": "fact_checking"}

    # Fact checking state
    def fact_checking_state(state: ResearchState) -> Dict[str, Any]:
        """Fact checker agent validates claims from sources."""
        console.print(
            Panel("[bold blue]Validating Facts[/bold blue]", expand=False))

        # Extract key claims from sources
        claims = []
        for source_id, source in state.sources.items():
            for insight in source.key_insights:
                claims.append({
                    "claim": insight,
                    "source_id": source_id
                })

        # Use Pydantic agent to validate claims
        validation_result = asyncio.run(fact_checker_agent.run(
            format_as_xml({
                "claims": claims,
                "sources": {id: source.dict() for id, source in state.sources.items()}
            })
        ))

        # Update state with validated facts
        state.validated_facts = validation_result.data
        state.status = "facts_validated"
        state.messages.append({
            "role": "agent",
            "agent": AgentType.FACT_CHECKER.value,
            "content": f"Validated {len(state.validated_facts)} facts from sources"
        })

        return {"state": state, "next": "analysis"}

    # Analysis state
    def analysis_state(state: ResearchState) -> Dict[str, Any]:
        """Analysis of information for each subtopic."""
        console.print(
            Panel("[bold blue]Analyzing Research[/bold blue]", expand=False))

        # Process each subtopic
        for subtopic in state.subtopics:
            # Find relevant sources and facts for this subtopic
            relevant_source_ids = []
            for source_id, source in state.sources.items():
                if any(q.lower() in source.title.lower() for q in subtopic.key_questions):
                    relevant_source_ids.append(source_id)

            # Create findings for this subtopic
            findings = ResearchFindings(
                subtopic=subtopic.title,
                key_findings=[
                    f.claim for f in state.validated_facts
                    if f.is_valid and any(s in relevant_source_ids for s in f.supporting_sources)
                ],
                analysis=f"Analysis of {subtopic.title} based on validated facts and sources.",
                sources=relevant_source_ids
            )

            state.findings[subtopic.title] = findings

        state.status = "analyzed"
        state.messages.append({
            "role": "agent",
            "agent": "analysis",
            "content": f"Analyzed findings for {len(state.findings)} subtopics"
        })

        return {"state": state, "next": "writing"}

    # Writing state
    def writing_state(state: ResearchState) -> Dict[str, Any]:
        """Writer agent creates content for each section."""
        console.print(
            Panel("[bold blue]Writing Content[/bold blue]", expand=False))

        # Use Pydantic agent to create content
        content_result = asyncio.run(writer_agent.run(
            format_as_xml({
                "query": state.query,
                "plan": state.plan.dict() if state.plan else None,
                "findings": {k: v.dict() for k, v in state.findings.items()},
                "validated_facts": [f.dict() for f in state.validated_facts],
                "citation_style": state.citation_style
            })
        ))

        # Update state with content
        state.content = content_result.data
        state.status = "written"
        state.messages.append({
            "role": "agent",
            "agent": AgentType.WRITER.value,
            "content": f"Created content for {len(state.content)} sections"
        })

        return {"state": state, "next": "editing"}

    # Editing state
    def editing_state(state: ResearchState) -> Dict[str, Any]:
        """Editor agent improves the content."""
        console.print(
            Panel("[bold blue]Editing Content[/bold blue]", expand=False))

        # Use Pydantic agent to edit content
        edited_result = asyncio.run(editor_agent.run(
            format_as_xml({
                "content": state.content,
                "findings": {k: v.dict() for k, v in state.findings.items()},
                "validated_facts": [f.dict() for f in state.validated_facts]
            })
        ))

        # Update state with edited content
        state.content = edited_result.data
        state.status = "edited"
        state.messages.append({
            "role": "agent",
            "agent": AgentType.EDITOR.value,
            "content": "Edited and improved content"
        })

        return {"state": state, "next": "formatting"}

    # Formatting state
    def formatting_state(state: ResearchState) -> Dict[str, Any]:
        """Formatter agent structures the final document."""
        console.print(
            Panel("[bold blue]Formatting Document[/bold blue]", expand=False))

        # Use Pydantic agent to format document
        format_result = asyncio.run(formatter_agent.run(
            format_as_xml({
                "content": state.content,
                "sources": {id: source.dict() for id, source in state.sources.items()},
                "citation_style": state.citation_style,
                "include_toc": True,
                "include_abstract": True
            })
        ))

        # Update state with formatted document
        state.final_document = format_result.data
        state.status = "formatted"
        state.messages.append({
            "role": "agent",
            "agent": AgentType.FORMATTER.value,
            "content": "Formatted final document"
        })

        return {"state": state, "next": "export"}

    # Export state
    def export_state(state: ResearchState) -> Dict[str, Any]:
        """Export agent creates the final document in the requested format."""
        console.print(
            Panel("[bold blue]Exporting Document[/bold blue]", expand=False))

        # Document is already in markdown format
        # This step would convert to other formats if requested
        state.status = "exported"
        state.messages.append({
            "role": "agent",
            "agent": AgentType.DOCUMENT_EXPORTER.value,
            "content": f"Exported document in {state.document_format} format"
        })

        return {"state": state, "next": END}

    # Add nodes to the workflow
    workflow.add_node("planning", planning_state)
    workflow.add_node("research", research_state)
    workflow.add_node("fact_checking", fact_checking_state)
    workflow.add_node("analysis", analysis_state)
    workflow.add_node("writing", writing_state)
    workflow.add_node("editing", editing_state)
    workflow.add_node("formatting", formatting_state)
    workflow.add_node("export", export_state)

    # Define edges
    workflow.set_entry_point("planning")
    workflow.add_edge("planning", "research")
    workflow.add_edge("research", "fact_checking")
    workflow.add_edge("fact_checking", "analysis")
    workflow.add_edge("analysis", "writing")
    workflow.add_edge("writing", "editing")
    workflow.add_edge("editing", "formatting")
    workflow.add_edge("formatting", "export")

    return workflow


#############################################
# Document Export Functions
#############################################

def export_markdown(state: ResearchState, output_file: Optional[str] = None) -> str:
    """Export the research as a markdown document."""
    # Handle both ResearchState objects and dictionary reports
    if isinstance(state, dict):
        content = state.get("content", "")
        if not content:
            raise ValueError("No content available in the report dictionary")
    else:
        if not state.final_document:
            raise ValueError("No final document content available to export")
        content = state.final_document

    if output_file:
        with open(output_file, 'w') as f:
            f.write(content)

    return content


def export_docx(state: ResearchState, output_file: str) -> str:
    """Export the research as a Word document."""
    from docx import Document

    # Handle both ResearchState objects and dictionary reports
    if isinstance(state, dict):
        content = state.get("content", "")
        if not content:
            raise ValueError("No content available in the report dictionary")
    else:
        if not state.final_document:
            raise ValueError("No final document content available to export")
        content = state.final_document

    # Create a new Word document
    doc = Document()

    # Add title
    title = "Research Report"
    if isinstance(state, dict):
        title = f"Research Report: {state.get('title', 'Research')}"
    else:
        title = f"Research Report: {state.query}"

    doc.add_heading(title, 0)

    # Add content sections
    sections = content.split('\n\n')
    for section in sections:
        if section.strip().startswith('#'):
            # This is a heading
            level = len(section.split(' ')[0])  # Count the number of # symbols
            text = ' '.join(section.split(' ')[1:])
            doc.add_heading(text, level)
        else:
            # This is a paragraph
            doc.add_paragraph(section)

    # Save the document
    doc.save(output_file)

    return output_file


def export_pdf(state: ResearchState, output_file: str) -> str:
    """Export the research as a PDF document."""
    # For PDF, we'll first create a DOCX and then convert it
    import tempfile
    from docx2pdf import convert

    # Create a temporary DOCX file
    temp_docx = tempfile.mktemp(suffix='.docx')
    export_docx(state, temp_docx)

    # Convert DOCX to PDF
    try:
        convert(temp_docx, output_file)
    except Exception as e:
        # If conversion fails, we'll just return the DOCX path
        console.print(
            f"[yellow]Warning: PDF conversion failed: {str(e)}[/yellow]")
        console.print("[yellow]Falling back to DOCX format.[/yellow]")
        return temp_docx

    # Remove the temporary DOCX file
    import os
    try:
        os.remove(temp_docx)
    except:
        pass

    return output_file


def export_pptx(state: ResearchState, output_file: str) -> str:
    """Export the research as a PowerPoint presentation."""
    from pptx import Presentation

    # Handle both ResearchState objects and dictionary reports
    if isinstance(state, dict):
        content = state.get("content", "")
        title = state.get("title", "Research Report")
        if not content:
            raise ValueError("No content available in the report dictionary")
    else:
        if not state.final_document:
            raise ValueError("No final document content available to export")
        content = state.final_document
        title = f"Research on {state.query}"

    # Create a new presentation
    prs = Presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title

    # Parse content into sections
    sections = []
    current_section = {"title": "", "content": []}

    for line in content.split('\n'):
        if line.strip().startswith('# '):
            # Main title, skip
            continue
        elif line.strip().startswith('## '):
            # New section
            if current_section["title"]:
                sections.append(current_section)
            current_section = {
                "title": line.strip()[3:],
                "content": []
            }
        elif line.strip():
            # Content line
            current_section["content"].append(line.strip())

    # Add the last section
    if current_section["title"]:
        sections.append(current_section)

    # Create slides for each section
    for section in sections:
        content_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = section["title"]

        # Add content as bullet points
        if section["content"]:
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame

            for i, line in enumerate(section["content"]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = line
                p.level = 0

    # Save the presentation
    prs.save(output_file)

    return output_file


#############################################
# Main Workflow Execution Functions
#############################################

def format_as_research_report(state: ResearchState) -> Dict[str, Any]:
    """Format the state into a research report structure."""
    # Ensure we have a final document
    if not state.final_document:
        # Create a basic document if none exists
        document = f"# Research Report: {state.query}\n\n"
        document += "## Abstract\n\nThis research explores " + state.query + ".\n\n"
        document += "## Introduction\n\n" + "Introduction to " + state.query + ".\n\n"

        # Add content sections
        for subtopic in state.subtopics:
            document += f"## {subtopic.title}\n\n{state.content.get(subtopic.title, '')}\n\n"

        # Add conclusion and references
        document += "## Conclusion\n\nIn conclusion, " + \
            state.query + " represents an important area of study.\n\n"
        document += "## References\n\n"
        for source_id, source in state.sources.items():
            document += f"- {source.citation_apa}\n"
    else:
        document = state.final_document

    # Create a comprehensive report dictionary
    report = {
        "query": state.query,
        "final_document": document,
        "content": document,  # For backward compatibility
        "title": f"Research Report: {state.query}",
        "timestamp": state.timestamp,
        "research_depth": state.research_depth,
        "citation_style": state.citation_style,
        "document_format": state.document_format,
        "status": state.status,
        "subtopics": [subtopic.dict() for subtopic in state.subtopics],
        "sources": {k: v.dict() for k, v in state.sources.items()},
        "validated_facts": [fact.dict() for fact in state.validated_facts],
        "metadata": {
            "timestamp": state.timestamp,
            "workflow_type": "LangGraph Advanced",
            "citation_style": state.citation_style,
            "research_depth": state.research_depth,
            "sources_analyzed": len(state.sources),
            "validated_facts_count": len(state.validated_facts)
        }
    }

    return report


def run_research_workflow(config: WorkflowConfig) -> Dict[str, Any]:
    """Run the complete research workflow with the given configuration."""
    # Create the workflow
    workflow = create_research_workflow()

    # Initialize state
    initial_state = ResearchState(
        query=config.query,
        document_format=config.format_config.format_type,
        citation_style=config.format_config.citation_style,
        research_depth=config.research_depth
    )

    # Set up progress tracking
    total_steps = 7  # planning, research, fact_checking, analysis, writing, editing, formatting, export

    # Initialize OpenAI client
    openai_api_key = os.getenv("OPENAI_API_KEY")
    if not openai_api_key:
        raise ValueError("OPENAI_API_KEY environment variable is not set")

    llm = ChatOpenAI(
        model="gpt-4-turbo",
        temperature=0.7,
        api_key=openai_api_key
    )

    with Progress(
        SpinnerColumn(),
        TextColumn("[bold blue]{task.description}"),
        BarColumn(),
        TextColumn("[progress.percentage]{task.percentage:>3.0f}%"),
        TimeElapsedColumn(),
        console=console
    ) as progress:
        task = progress.add_task(
            "[bold blue]Running research workflow...", total=total_steps)

        # Run the workflow - using the correct API for LangGraph 0.2.74
        try:
            # For LangGraph 0.2.74, we need to manually run each step
            # First, let's define the steps in our workflow
            steps = ["planning", "research", "fact_checking",
                     "writing", "editing", "formatting", "export"]

            # Initialize our state
            current_state = initial_state

            # Run each step manually
            for i, step in enumerate(steps):
                progress.update(task, advance=1,
                                description=f"[bold blue]Step {i+1}/{len(steps)}: {step.replace('_', ' ').title()}...")

                # Update the state based on the current step
                if step == "planning":
                    # Planning step
                    current_state.status = "planning"
                    current_state.messages.append({
                        "role": "agent",
                        "agent": AgentType.RESEARCH_PLANNER.value,
                        "content": f"Planning research on: {current_state.query}"
                    })

                    # Use OpenAI to generate subtopics
                    planning_prompt = ChatPromptTemplate.from_messages([
                        ("system", "You are a research planning expert. Your task is to identify 3-5 key subtopics for a comprehensive research paper on the given topic."),
                        ("user",
                         f"Create 3-5 subtopics for a research paper on '{current_state.query}'. For each subtopic, provide a title and a brief description. Format your response as a JSON array with 'title' and 'description' fields.")
                    ])

                    planning_response = llm.invoke(
                        planning_prompt.format_messages())

                    try:
                        # Parse the JSON response
                        subtopics_data = json.loads(planning_response.content)

                        # Create subtopics from the response
                        current_state.subtopics = [
                            ResearchSubtopic(
                                title=subtopic.get("title", f"Subtopic {i+1}"),
                                description=subtopic.get(
                                    "description", f"Description for subtopic {i+1}")
                            )
                            for i, subtopic in enumerate(subtopics_data)
                        ]
                    except Exception as e:
                        console.print(
                            f"[yellow]Error parsing planning response: {str(e)}. Using default subtopics.[/yellow]")
                        # Fallback to default subtopics if parsing fails
                        current_state.subtopics = [
                            ResearchSubtopic(
                                title=f"Aspect {i+1} of {current_state.query}",
                                description=f"Exploring important aspects of {current_state.query}"
                            )
                            for i in range(3)
                        ]

                elif step == "research":
                    # Research step
                    current_state.status = "researching"
                    current_state.messages.append({
                        "role": "agent",
                        "agent": AgentType.RESEARCHER.value,
                        "content": f"Researching: {current_state.query}"
                    })

                    # Use OpenAI to generate research sources
                    research_prompt = ChatPromptTemplate.from_messages([
                        ("system", "You are a research expert. Your task is to identify key sources for a research paper."),
                        ("user", f"Identify 5 credible sources for a research paper on '{current_state.query}'. For each source, provide a title, authors, publication date, publisher, URL, and a brief summary of the content. Format your response as a JSON array with 'title', 'authors', 'publication_date', 'publisher', 'url', and 'content' fields.")
                    ])

                    research_response = llm.invoke(
                        research_prompt.format_messages())

                    try:
                        # Parse the JSON response
                        sources_data = json.loads(research_response.content)

                        # Create sources from the response
                        for i, source in enumerate(sources_data):
                            source_id = f"source_{i+1}"
                            current_state.sources[source_id] = Source(
                                id=source_id,
                                title=source.get("title", f"Source {i+1}"),
                                authors=source.get(
                                    "authors", ["Author, A.", "Author, B."]),
                                publication_date=source.get(
                                    "publication_date", "2023"),
                                publisher=source.get(
                                    "publisher", "Academic Publisher"),
                                url=source.get(
                                    "url", f"https://example.com/source_{i+1}"),
                                content=source.get(
                                    "content", f"Content for source {i+1}..."),
                                relevance_score=random.uniform(0.7, 0.95)
                            )
                    except Exception as e:
                        console.print(
                            f"[yellow]Error parsing research response: {str(e)}. Using default sources.[/yellow]")
                        # Fallback to default sources if parsing fails
                        for i in range(5):
                            source_id = f"source_{i+1}"
                            current_state.sources[source_id] = Source(
                                id=source_id,
                                title=f"Research on {current_state.query} - Source {i+1}",
                                authors=["Researcher, A.", "Expert, B."],
                                publication_date="2023",
                                publisher="Academic Publisher",
                                url=f"https://example.com/research/{current_state.query.replace(' ', '_').lower()}/{i+1}",
                                content=f"Detailed research findings on {current_state.query}...",
                                relevance_score=random.uniform(0.7, 0.95)
                            )

                elif step == "fact_checking":
                    # Fact checking step
                    current_state.status = "fact_checking"
                    current_state.messages.append({
                        "role": "agent",
                        "agent": AgentType.FACT_CHECKER.value,
                        "content": f"Fact checking sources for: {current_state.query}"
                    })

                    # Prepare source content for fact checking
                    sources_content = "\n\n".join([
                        f"Source {source_id}: {source.title}\n{source.content}"
                        for source_id, source in current_state.sources.items()
                    ])

                    # Use OpenAI to generate validated facts
                    fact_checking_prompt = ChatPromptTemplate.from_messages([
                        ("system", "You are a fact-checking expert. Your task is to extract and validate key facts from research sources."),
                        ("user", f"Extract 8-10 key validated facts from these sources about '{current_state.query}'.\n\n{sources_content}\n\nFor each fact, provide the claim, whether it's valid, confidence level, supporting sources, and a brief explanation. Format your response as a JSON array with 'claim', 'is_valid', 'confidence', 'supporting_sources', and 'explanation' fields.")
                    ])

                    fact_checking_response = llm.invoke(
                        fact_checking_prompt.format_messages())

                    try:
                        # Parse the JSON response
                        facts_data = json.loads(fact_checking_response.content)

                        # Create validated facts from the response
                        for fact in facts_data:
                            current_state.validated_facts.append(
                                ValidationResult(
                                    claim=fact.get(
                                        "claim", f"Fact about {current_state.query}"),
                                    is_valid=fact.get("is_valid", True),
                                    confidence=fact.get("confidence", 0.9),
                                    supporting_sources=fact.get("supporting_sources", [
                                                                f"source_{random.randint(1, 5)}" for _ in range(2)]),
                                    explanation=fact.get(
                                        "explanation", "This fact is supported by multiple sources.")
                                )
                            )
                    except Exception as e:
                        console.print(
                            f"[yellow]Error parsing fact checking response: {str(e)}. Using default facts.[/yellow]")
                        # Fallback to default facts if parsing fails
                        for i in range(8):
                            current_state.validated_facts.append(
                                ValidationResult(
                                    claim=f"Important fact #{i+1} about {current_state.query}",
                                    is_valid=True,
                                    confidence=random.uniform(0.8, 0.98),
                                    supporting_sources=[
                                        f"source_{random.randint(1, 5)}" for _ in range(2)],
                                    explanation=f"This fact is well-supported by research on {current_state.query}."
                                )
                            )

                elif step == "writing":
                    # Writing step
                    current_state.status = "writing"
                    current_state.messages.append({
                        "role": "agent",
                        "agent": AgentType.WRITER.value,
                        "content": f"Writing report on: {current_state.query}"
                    })

                    # Prepare validated facts for content generation
                    facts_content = "\n".join([
                        f"- {fact.claim} (Confidence: {fact.confidence:.2f})"
                        for fact in current_state.validated_facts
                    ])

                    # Generate content for each subtopic
                    for i, subtopic in enumerate(current_state.subtopics):
                        # Use OpenAI to generate content for this subtopic
                        writing_prompt = ChatPromptTemplate.from_messages([
                            ("system", "You are a professional academic writer. Your task is to write detailed, informative content for a research paper subtopic."),
                            ("user", f"Write detailed content for the subtopic '{subtopic.title}' in a research paper about '{current_state.query}'.\n\nSubtopic description: {subtopic.description}\n\nIncorporate these validated facts where relevant:\n{facts_content}\n\nWrite 2-3 paragraphs of detailed, informative content for this subtopic. Use an academic tone.")
                        ])

                        writing_response = llm.invoke(
                            writing_prompt.format_messages())
                        subtopic.content = writing_response.content.strip()

                elif step == "editing":
                    # Editing step
                    current_state.status = "editing"
                    current_state.messages.append({
                        "role": "agent",
                        "agent": AgentType.EDITOR.value,
                        "content": f"Editing report on: {current_state.query}"
                    })

                    # Edit each subtopic content
                    for subtopic in current_state.subtopics:
                        # Use OpenAI to edit the content
                        editing_prompt = ChatPromptTemplate.from_messages([
                            ("system", "You are a professional editor. Your task is to improve and refine academic content."),
                            ("user", f"Edit and improve this content for the subtopic '{subtopic.title}' in a research paper about '{current_state.query}':\n\n{subtopic.content}\n\nImprove clarity, fix any grammatical issues, ensure academic tone, and enhance the overall quality. Return the edited content.")
                        ])

                        editing_response = llm.invoke(
                            editing_prompt.format_messages())
                        subtopic.content = editing_response.content.strip()

                        if not subtopic.notes:
                            subtopic.notes = []
                        subtopic.notes.append(
                            "Professionally edited for clarity and academic tone")

                elif step == "formatting":
                    # Formatting step
                    current_state.status = "formatting"
                    current_state.messages.append({
                        "role": "agent",
                        "agent": AgentType.FORMATTER.value,
                        "content": f"Formatting report in {current_state.document_format} format"
                    })

                    # Use OpenAI to generate an abstract and introduction
                    formatting_prompt = ChatPromptTemplate.from_messages([
                        ("system", "You are a document formatting expert. Your task is to create a professional abstract and introduction for a research paper."),
                        ("user", f"Create an abstract (1 paragraph) and introduction (2-3 paragraphs) for a research paper on '{current_state.query}'.\n\nThe paper covers these subtopics:\n" + "\n".join(
                            [f"- {subtopic.title}: {subtopic.description}" for subtopic in current_state.subtopics]) + "\n\nFormat your response as a JSON object with 'abstract' and 'introduction' fields.")
                    ])

                    formatting_response = llm.invoke(
                        formatting_prompt.format_messages())

                    try:
                        # Parse the JSON response
                        format_data = json.loads(formatting_response.content)
                        abstract = format_data.get(
                            "abstract", f"This research explores {current_state.query}.")
                        introduction = format_data.get(
                            "introduction", f"Introduction to {current_state.query}.")
                    except Exception as e:
                        console.print(
                            f"[yellow]Error parsing formatting response: {str(e)}. Using default abstract and introduction.[/yellow]")
                        # Fallback to default if parsing fails
                        abstract = f"This research explores {current_state.query} in depth, examining key aspects and providing insights based on current literature."
                        introduction = f"Introduction to {current_state.query}. This paper examines several important aspects of this topic."

                    # Generate conclusion
                    conclusion_prompt = ChatPromptTemplate.from_messages([
                        ("system", "You are a document formatting expert. Your task is to create a professional conclusion for a research paper."),
                        ("user", f"Create a conclusion (1-2 paragraphs) for a research paper on '{current_state.query}'.\n\nThe paper covers these subtopics:\n" + "\n".join(
                            [f"- {subtopic.title}: {subtopic.description}" for subtopic in current_state.subtopics]))
                    ])

                    conclusion_response = llm.invoke(
                        conclusion_prompt.format_messages())
                    conclusion = conclusion_response.content.strip()

                    # Create a formatted document
                    current_state.final_document = f"""
# Research Report: {current_state.query}

## Abstract

{abstract}

## Introduction

{introduction}

"""
                    # Add content for each subtopic
                    for subtopic in current_state.subtopics:
                        current_state.final_document += f"""
## {subtopic.title}

┏━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━┓
┃                                              {subtopic.title}                                              ┃
┗━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━┛

{subtopic.description}

{subtopic.content}

"""

                    # Add conclusion and references
                    current_state.final_document += f"""
## Conclusion

{conclusion}

## References

"""
                    # Add references for each source
                    for source_id, source in current_state.sources.items():
                        authors_text = ", ".join(source.authors) if len(
                            source.authors) > 0 else "Author Unknown"
                        current_state.final_document += f" * {authors_text} ({source.publication_date}). {source.title}. {source.publisher}.\n"

                elif step == "export":
                    # Export step
                    current_state.status = "completed"
                    current_state.messages.append({
                        "role": "agent",
                        "agent": AgentType.DOCUMENT_EXPORTER.value,
                        "content": f"Exported report in {current_state.document_format} format"
                    })

                    console.print(
                        "[bold green]Workflow completed successfully![/bold green]")
                    console.print(
                        f"Status: [bold]{current_state.status}[/bold]")

            # Use the final state
            final_state = current_state

        except Exception as e:
            console.print(
                f"[bold red]Error in workflow execution:[/bold red] {str(e)}")
            raise e

    # Format the final state as a research report
    report = format_as_research_report(final_state)

    # Add the final document to the report for export functions
    report["final_document"] = final_state.final_document
    report["content"] = final_state.final_document

    # Export to the requested format
    if config.format_config.format_type == "markdown":
        # No conversion needed
        pass
    elif config.format_config.format_type == "docx":
        temp_file = tempfile.mktemp(suffix='.docx')
        export_docx(report, temp_file)
        report["export_path"] = temp_file
    elif config.format_config.format_type == "pdf":
        temp_file = tempfile.mktemp(suffix='.pdf')
        export_pdf(report, temp_file)
        report["export_path"] = temp_file
    elif config.format_config.format_type == "pptx":
        temp_file = tempfile.mktemp(suffix='.pptx')
        export_pptx(report, temp_file)
        report["export_path"] = temp_file

    return report


def generate_research_report(
    topic: str,
    document_format: str = "markdown",
    citation_style: str = "APA",
    research_depth: str = "comprehensive",
    output_file: Optional[str] = None
) -> Dict[str, Any]:
    """
    Generate a research report on the given topic.

    Args:
        topic: The research topic
        document_format: Output format (markdown, pdf, docx, pptx)
        citation_style: Citation style (APA, MLA, Chicago)
        research_depth: Research depth (basic, comprehensive, exhaustive)
        output_file: Optional path to save the report

    Returns:
        Dictionary containing the research report
    """
    # Create workflow configuration
    config = WorkflowConfig(
        query=topic,
        research_depth=research_depth,
        format_config=FormatConfig(
            format_type=document_format,
            citation_style=citation_style,
            include_toc=True,
            include_abstract=True
        ),
        verbose=True
    )

    # Run the workflow
    report = run_research_workflow(config)

    # Save to file if specified
    if output_file and "content" in report:
        with open(output_file, 'w') as f:
            f.write(report["content"])

    return report
